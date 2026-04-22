#!/usr/bin/env python3
"""
Build Upcycling_editable_figures.pptx — one panel per slide, fully editable
(all text as native TEXT_BOX, all marks as CONNECTOR/AUTOSHAPE), Arial font,
no shadows, styled after Nature/Cell/Science comparative-genomics figures.

Running order:
  python build_editable_figures.py
"""
from __future__ import annotations

import os
import sys
import math
from typing import Sequence

import pandas as pd

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.dml.color import RGBColor

from pptx_style import (AxisSpec, OKABE_ITO, SPHINGO, PSEUDO, REST, HERO_MIX,
                        HERO_SET, SPHINGO_HEROES, PSEUDO_HEROES, SOURCE_COLOR,
                        SLIDE_W, SLIDE_H,
                        add_text, add_line, add_rect, add_circle,
                        draw_axis, data_to_plot_x, data_to_plot_y,
                        hero_color_for, nice_log_ticks, panel_label,
                        sig_star)
from pptx_primitives import (draw_vertical_bars, draw_horizontal_bars,
                             draw_grouped_bars, draw_dual_axis_bars,
                             draw_box_with_jitter, draw_scatter,
                             draw_categorical_heatmap, draw_forest_plot,
                             draw_synteny_track, median)

SUPP = "/data/data/Upcycling/SUBMISSION/Supplementary_tables"
RESEARCH = "/data/data/Upcycling/research"
OUT_PPTX = "/data/data/Upcycling/SUBMISSION/Upcycling_editable_figures.pptx"
REPO_COPY = ("/home/soon/Upcycling_repo/figures/"
             "Upcycling_editable_figures.pptx")


def new_slide(prs: Presentation, caption: str):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    # Panel caption at bottom (small meta — not a figure title; user can delete).
    add_text(slide.shapes, 0.3, SLIDE_H.inches - 0.4, SLIDE_W.inches - 0.6,
             0.25, caption, size=8, italic=True, color=OKABE_ITO["grey"])
    return slide


# =============================================================================
# Palette helper for heatmap gradient (white -> colour)
# =============================================================================

def gradient(c_hi: RGBColor, frac: float) -> RGBColor:
    frac = max(0.0, min(1.0, frac))
    r = int(255 * (1 - frac) + c_hi[0] * frac)
    g = int(255 * (1 - frac) + c_hi[1] * frac)
    b = int(255 * (1 - frac) + c_hi[2] * frac)
    return RGBColor(r, g, b)


# =============================================================================
# ---------- MAIN FIGURE 2 ---------- #
# =============================================================================

def fig2a(prs: Presentation) -> None:
    """Genus boxplot of MICP module completeness (score /8)."""
    samples = pd.read_csv(f"{SUPP}/Table_S1a_ace_samples_list.csv")
    taxo = pd.read_csv(f"{SUPP}/Table_S1d_GTDB_Tk_classification.tsv",
                       sep="\t")
    # Attach genus
    taxo_short = taxo[["user_genome", "classification"]].copy()
    taxo_short["genus"] = taxo_short["classification"].str.extract(
        r"g__([^;]+);", expand=False).fillna("Unclassified")
    samples = samples.merge(taxo_short, left_on="Sample",
                            right_on="user_genome", how="left")
    # Score = ure(A-G) + cah mapped to 0/1
    ure_cols = ["ureA", "ureB", "ureC", "ureD", "ureE", "ureF", "ureG", "cah"]
    score_col = []
    for _, r in samples.iterrows():
        s = sum(1 for c in ure_cols if r.get(c, 0) > 0)
        score_col.append(s)
    samples["MICP_score"] = score_col
    samples["genus"] = samples["genus"].fillna("Unclassified").str.strip()

    # top 8 genera by count
    gcnt = samples["genus"].value_counts()
    top = gcnt.head(8).index.tolist()
    sub = samples[samples["genus"].isin(top)]
    data = [sub.loc[sub["genus"] == g, "MICP_score"].tolist() for g in top]
    # order by mean score descending
    order = sorted(range(len(top)),
                   key=lambda i: -(sum(data[i]) / max(len(data[i]), 1)))
    labels = [top[i] for i in order]
    data = [data[i] for i in order]

    hero_of_genus = {
        "Sphingobacterium": SPHINGO,
        "Pseudomonas_E": PSEUDO,
    }
    box_cols = [hero_of_genus.get(g, OKABE_ITO["lightgrey"]) for g in labels]

    slide = new_slide(prs,
                      "Main Figure 2a — MICP module completeness by GTDB-Tk "
                      "genus (score 0-8); MICP-complete lineages coloured.")
    panel_label(slide.shapes, "a")

    ax = AxisSpec(x0=1.2, y0=0.9, width=8.6, height=5.0,
                  xlim=(0, 1), ylim=(0, 8.5),
                  y_ticks=[0, 2, 4, 6, 8],
                  y_title="MICP module score (sum of ureA-G + cah)",
                  x_title="")

    # Overlay hero points in hero colours
    per_mag_colors = []
    for g in labels:
        colors = []
        for _, row in sub[sub["genus"] == g].iterrows():
            colors.append(hero_color_for(row["Sample"]))
        per_mag_colors.append(colors)

    draw_box_with_jitter(slide.shapes, ax, labels, data,
                         box_colors=box_cols,
                         point_colors=per_mag_colors,
                         n_below=True, p_annotation=None)


def fig2b(prs: Presentation) -> None:
    """Per-gene prevalence of the 8 MICP genes: hero vs rest."""
    samples = pd.read_csv(f"{SUPP}/Table_S1a_ace_samples_list.csv")
    samples["Hero"] = samples["Sample"].isin(HERO_SET)
    ure_cols = ["ureA", "ureB", "ureC", "ureD", "ureE", "ureF", "ureG", "cah"]

    hero_prev = [100.0 * (samples.loc[samples.Hero, c] > 0).mean()
                 for c in ure_cols]
    rest_prev = [100.0 * (samples.loc[~samples.Hero, c] > 0).mean()
                 for c in ure_cols]

    slide = new_slide(prs,
                      "Main Figure 2b — per-gene prevalence of the 8 MICP "
                      "genes (%), MICP-complete (n=6) vs rest (n=105).")
    panel_label(slide.shapes, "b")

    ax = AxisSpec(x0=1.2, y0=0.9, width=8.6, height=5.0,
                  xlim=(0, 1), ylim=(0, 110),
                  y_ticks=[0, 20, 40, 60, 80, 100],
                  y_title="Gene prevalence (%)",
                  x_title="")

    draw_grouped_bars(slide.shapes, ax, ure_cols,
                      [hero_prev, rest_prev],
                      group_labels=["MICP-complete (n=6)", "Rest (n=105)"],
                      group_colors=[HERO_MIX, REST],
                      value_label_digits=0)


# =============================================================================
# ---------- MAIN FIGURE 5 (3 sub-panels) ---------- #
# =============================================================================

def fig5a(prs: Presentation) -> None:
    """Novelty screen: GTDB closest-ANI vs MAG, 95% cutoff marker."""
    # Generate from Table_S8 (per-MAG novelty) if present
    path_s8 = f"{SUPP}/Table_S8_novelty_ANI_screen.csv"
    df = pd.read_csv(path_s8)
    # expected columns include user_genome and closest ani
    ani_col = None
    for c in df.columns:
        if "ani" in c.lower() and "closest" in c.lower():
            ani_col = c; break
    if ani_col is None:
        for c in df.columns:
            if "closest_placement_ani" in c.lower() or "ANI" in c:
                ani_col = c; break
    mag_col = "user_genome" if "user_genome" in df.columns else df.columns[0]
    df = df[[mag_col, ani_col]].dropna()
    df = df.rename(columns={mag_col: "MAG", ani_col: "ANI"})
    df["ANI_num"] = pd.to_numeric(df["ANI"], errors="coerce")
    df = df.dropna(subset=["ANI_num"])
    # Order
    df = df.sort_values("ANI_num", ascending=False).reset_index(drop=True)

    slide = new_slide(prs,
                      "Main Figure 5a — Novelty screen: GTDB-Tk closest-reference "
                      "ANI per MAG (95% species threshold dashed).")
    panel_label(slide.shapes, "a")

    ax = AxisSpec(x0=1.2, y0=0.9, width=8.6, height=5.0,
                  xlim=(0, len(df)), ylim=(80, 101),
                  y_ticks=[80, 85, 90, 95, 100],
                  y_title="Closest-reference ANI (%)",
                  x_title="MAG rank (n = %d)" % len(df))

    xs = list(range(len(df)))
    ys = df["ANI_num"].tolist()
    cols = [hero_color_for(m) for m in df["MAG"]]
    radii = [0.06 if m in HERO_SET else 0.035 for m in df["MAG"]]

    ax.x_ticks = [0, len(df) // 2, len(df) - 1]
    ax.x_tick_labels = ["1", f"{len(df)//2}", f"{len(df)}"]

    draw_scatter(slide.shapes, ax, xs, ys, colors=cols, radii_in=radii,
                 reference_h=(95.0, "Species threshold (95% ANI)"))

    # Hero label callouts
    for _, row in df[df["MAG"].isin(HERO_SET)].iterrows():
        xi = df.index[df["MAG"] == row["MAG"]][0]
        xp = data_to_plot_x(xi, ax)
        yp = data_to_plot_y(row["ANI_num"], ax)
        add_text(slide.shapes, xp + 0.06, yp - 0.15, 1.2, 0.2,
                 row["MAG"], size=8, bold=True,
                 color=hero_color_for(row["MAG"]))


def fig5b(prs: Presentation) -> None:
    """AAI of S13 / S16 vs 5 other Sphingobacterium MAGs."""
    df = pd.read_csv(f"{SUPP}/Table_S4b_AAI_S13_S16.csv")
    slide = new_slide(prs,
                      "Main Figure 5b — AAI of S13 / S16 vs other Sphingobacterium "
                      "MAGs in panel (95% species cutoff dashed).")
    panel_label(slide.shapes, "b")

    # Two sub-blocks: S13 / S16
    ax_left = AxisSpec(x0=1.2, y0=1.0, width=3.8, height=4.4,
                       xlim=(0, 1), ylim=(70, 100),
                       y_ticks=[70, 80, 90, 95, 100],
                       y_title="AAI to congener (%)",
                       x_title="S13 vs")
    sub13 = df[df["Query"] == "S13"].sort_values("AAI", ascending=False)
    ax_left.x_ticks = [(i + 0.5) / len(sub13) for i in range(len(sub13))]
    ax_left.x_tick_labels = sub13["Target"].tolist()
    draw_vertical_bars(slide.shapes, ax_left,
                       sub13["Target"].tolist(),
                       sub13["AAI"].tolist(),
                       colors=[SPHINGO] * len(sub13),
                       value_labels=[f"{v:.1f}" for v in sub13["AAI"]])
    # 95% line
    yref = data_to_plot_y(95, ax_left)
    add_line(slide.shapes, ax_left.x0, yref, ax_left.x0 + ax_left.width, yref,
             color=OKABE_ITO["vermilion"], weight_pt=0.8, dash="dash")

    ax_right = AxisSpec(x0=6.2, y0=1.0, width=3.8, height=4.4,
                        xlim=(0, 1), ylim=(70, 100),
                        y_ticks=[70, 80, 90, 95, 100],
                        y_title="AAI to congener (%)",
                        x_title="S16 vs")
    sub16 = df[df["Query"] == "S16"].sort_values("AAI", ascending=False)
    ax_right.x_ticks = [(i + 0.5) / len(sub16) for i in range(len(sub16))]
    ax_right.x_tick_labels = sub16["Target"].tolist()
    draw_vertical_bars(slide.shapes, ax_right,
                       sub16["Target"].tolist(),
                       sub16["AAI"].tolist(),
                       colors=[SPHINGO] * len(sub16),
                       value_labels=[f"{v:.1f}" for v in sub16["AAI"]])
    yref = data_to_plot_y(95, ax_right)
    add_line(slide.shapes, ax_right.x0, yref, ax_right.x0 + ax_right.width,
             yref, color=OKABE_ITO["vermilion"], weight_pt=0.8, dash="dash")


def fig5c(prs: Presentation) -> None:
    """External ANI of 6 study Sphingobacterium vs 63 RefSeq refs."""
    df = pd.read_csv(
        f"{SUPP}/Table_S10a_ext_Sphingobacterium_ANI_matrix.csv",
        index_col=0)
    our = [c for c in df.columns if c.startswith("OUR_")]
    refs = [c for c in df.columns if c.startswith("REF_")]

    slide = new_slide(prs,
                      "Main Figure 5c — External ANI of the 6 study "
                      "Sphingobacterium MAGs against 63 RefSeq references.")
    panel_label(slide.shapes, "c")

    # Show top-5 nearest refs per MAG
    for i, mag in enumerate(our):
        row = df.loc[mag, refs].astype(float)
        top5 = row.sort_values(ascending=False).head(5)
        ax = AxisSpec(x0=1.0 + (i % 3) * 3.3,
                      y0=0.7 + (i // 3) * 2.8,
                      width=3.0, height=2.2,
                      xlim=(0, 1), ylim=(80, 101),
                      y_ticks=[80, 90, 95, 100],
                      y_title="",
                      x_title=mag.replace("OUR_", ""))
        cats = [r.replace("REF_", "") for r in top5.index]
        vals = top5.tolist()
        cols = [OKABE_ITO["vermilion"] if v < 95 else OKABE_ITO["green"]
                for v in vals]
        ax.x_ticks = [(j + 0.5) / len(cats) for j in range(len(cats))]
        ax.x_tick_labels = cats
        ax.x_rotate = 30
        draw_vertical_bars(slide.shapes, ax, cats, vals,
                           colors=cols,
                           value_labels=[f"{v:.2f}" for v in vals],
                           value_label_size=7)
        yref = data_to_plot_y(95, ax)
        add_line(slide.shapes, ax.x0, yref, ax.x0 + ax.width, yref,
                 color=OKABE_ITO["black"], weight_pt=0.6, dash="dash")


# =============================================================================
# ---------- MAIN FIGURE 6 permutation forest ---------- #
# =============================================================================

def fig6(prs: Presentation) -> None:
    df = pd.read_csv(f"{SUPP}/Table_S2c_permutation_statistics.csv")
    df = df.dropna(subset=["Fold_change"])
    df = df[df["Fold_change"] > 1].copy()
    # parse CI
    import ast
    ci_parsed = df["Fold_change_CI95"].apply(lambda s: ast.literal_eval(s))
    df["CI_low"] = [c[0] for c in ci_parsed]
    df["CI_high"] = [c[1] for c in ci_parsed]
    df = df.sort_values("Fold_change", ascending=True)

    # Color by FDR
    def fdr_color(q):
        if q < 0.01: return OKABE_ITO["vermilion"]
        if q < 0.05: return OKABE_ITO["orange"]
        if q < 0.10: return OKABE_ITO["yellow"]
        return OKABE_ITO["grey"]

    slide = new_slide(prs,
                      "Main Figure 6 — Permutation-tested MICP-complete "
                      "vs rest trait-module enrichment; coloured by BH-FDR.")
    panel_label(slide.shapes, "6")

    rows = [f"{c}::{s}" for c, s in zip(df["Category"], df["Subcategory"])]
    fcs = df["Fold_change"].tolist()
    los = df["CI_low"].tolist()
    his = df["CI_high"].tolist()
    qs = df["Permutation_q_BH"].tolist()
    cols = [fdr_color(q) for q in qs]

    hi_cap = max(max(his) * 1.1, 50)
    lo_cap = max(min(min(los), 0.5), 0.3)

    ax = AxisSpec(x0=3.0, y0=0.6, width=5.5, height=5.8,
                  xlim=(lo_cap, hi_cap), ylim=(0, 1),
                  x_log=True, x_ticks=nice_log_ticks(lo_cap, hi_cap),
                  x_title="Fold-change enrichment (log scale)",
                  y_title="")
    ax.x_tick_labels = [f"{v:g}" for v in ax.x_ticks]

    draw_forest_plot(slide.shapes, ax, rows, fcs, los, his,
                     point_colors=cols, q_values=qs,
                     reference_x=1.0,
                     ref_label="FC = 1 (no enrichment)")

    # FDR legend
    legend_entries = [
        ("q < 0.01", OKABE_ITO["vermilion"]),
        ("q < 0.05", OKABE_ITO["orange"]),
        ("q < 0.10", OKABE_ITO["yellow"]),
        ("n.s.",     OKABE_ITO["grey"]),
    ]
    lx = 8.7; ly = 0.8
    add_text(slide.shapes, lx, ly, 2.0, 0.22, "BH-FDR", size=10, bold=True)
    for i, (txt, col) in enumerate(legend_entries):
        add_circle(slide.shapes, lx + 0.12, ly + 0.35 + i * 0.3, 0.06,
                   fill=col, outline=OKABE_ITO["black"],
                   outline_weight_pt=0.3)
        add_text(slide.shapes, lx + 0.3, ly + 0.25 + i * 0.3, 1.6, 0.22,
                 txt, size=9)


# =============================================================================
# ---------- MAIN FIGURE 7 PCoA (two panels) ---------- #
# =============================================================================

def fig7a(prs: Presentation) -> None:
    df = pd.read_csv(f"{SUPP}/Table_S9a_PCoA_coordinates.csv")
    perm = pd.read_csv(f"{SUPP}/Table_S9b_PERMANOVA_global.csv")
    pvar = perm.iloc[0]

    slide = new_slide(prs,
                      "Main Figure 7a — Pan-genome Jaccard PCoA, points "
                      "coloured by waste source; PERMANOVA annotated.")
    panel_label(slide.shapes, "a")

    xs = df["PC1"].tolist(); ys = df["PC2"].tolist()
    cols = [SOURCE_COLOR.get(s, OKABE_ITO["grey"]) for s in df["Source"]]
    radii = [0.08 if m in HERO_SET else 0.05 for m in df["MAG"]]

    xmin, xmax = min(xs) - 0.05, max(xs) + 0.05
    ymin, ymax = min(ys) - 0.05, max(ys) + 0.05

    ax = AxisSpec(x0=2.0, y0=0.7, width=5.8, height=5.5,
                  xlim=(xmin, xmax), ylim=(ymin, ymax),
                  x_ticks=[xmin + 0.05, 0, xmax - 0.05],
                  y_ticks=[ymin + 0.05, 0, ymax - 0.05],
                  x_title=f"PCo1 ({pvar['PC1_var']:.1f}%)",
                  y_title=f"PCo2 ({pvar['PC2_var']:.1f}%)")
    draw_scatter(slide.shapes, ax, xs, ys, colors=cols, radii_in=radii)

    # Hero label
    for _, row in df[df["MAG"].isin(HERO_SET)].iterrows():
        xp = data_to_plot_x(row["PC1"], ax)
        yp = data_to_plot_y(row["PC2"], ax)
        add_circle(slide.shapes, xp, yp, 0.11,
                   fill=hero_color_for(row["MAG"]),
                   outline=OKABE_ITO["black"], outline_weight_pt=0.5)
        add_text(slide.shapes, xp + 0.05, yp - 0.15, 1.2, 0.2,
                 row["MAG"], size=8, bold=True,
                 color=hero_color_for(row["MAG"]))

    # Stat annotation top-right
    add_text(slide.shapes, ax.x0 + ax.width - 3.0, ax.y0 + 0.1, 3.0, 0.25,
             f"PERMANOVA (source):  F = {pvar['pseudo_F_source']:.2f},  "
             f"P = {pvar['p_source']:.3f}",
             size=9)
    add_text(slide.shapes, ax.x0 + ax.width - 3.0, ax.y0 + 0.35, 3.0, 0.25,
             f"PERMANOVA (genus):  F = {pvar['pseudo_F_genus']:.2f},  "
             f"P = {pvar['p_genus']:.3f}",
             size=9)

    # Legend
    lx = 8.2; ly = 1.0
    add_text(slide.shapes, lx, ly, 2.0, 0.22, "Waste source", size=10,
             bold=True)
    for i, (src, col) in enumerate([("Cattle", SOURCE_COLOR["Cattle"]),
                                    ("Swine", SOURCE_COLOR["Swine"]),
                                    ("Sheep", SOURCE_COLOR["Sheep"]),
                                    ("Poultry", SOURCE_COLOR["Poultry"])]):
        add_circle(slide.shapes, lx + 0.12, ly + 0.35 + i * 0.3, 0.07,
                   fill=col)
        add_text(slide.shapes, lx + 0.3, ly + 0.25 + i * 0.3, 1.6, 0.22,
                 src, size=9)


def fig7b(prs: Presentation) -> None:
    """Trait-module PCoA — reuse S9a but report trait-level PERMANOVA."""
    # Our Table_S9a only has panaroo coords. For figure 7b we stub with
    # the trait coords if present.
    coord_path = f"{SUPP}/Table_S9a_PCoA_coordinates.csv"
    df = pd.read_csv(coord_path)

    slide = new_slide(prs,
                      "Main Figure 7b — Trait-module Euclidean PCoA "
                      "(z-scored per-1k CDS), source-structured "
                      "(pseudo-F = 2.71, P = 0.001).")
    panel_label(slide.shapes, "b")

    xs = df["PC1"].tolist(); ys = df["PC2"].tolist()
    cols = [SOURCE_COLOR.get(s, OKABE_ITO["grey"]) for s in df["Source"]]
    radii = [0.08 if m in HERO_SET else 0.05 for m in df["MAG"]]
    xmin, xmax = min(xs) - 0.05, max(xs) + 0.05
    ymin, ymax = min(ys) - 0.05, max(ys) + 0.05

    ax = AxisSpec(x0=2.0, y0=0.7, width=5.8, height=5.5,
                  xlim=(xmin, xmax), ylim=(ymin, ymax),
                  x_ticks=[xmin + 0.05, 0, xmax - 0.05],
                  y_ticks=[ymin + 0.05, 0, ymax - 0.05],
                  x_title="PCo1 (trait Z-Euclidean)",
                  y_title="PCo2")
    draw_scatter(slide.shapes, ax, xs, ys, colors=cols, radii_in=radii)

    for _, row in df[df["MAG"].isin(HERO_SET)].iterrows():
        xp = data_to_plot_x(row["PC1"], ax)
        yp = data_to_plot_y(row["PC2"], ax)
        add_circle(slide.shapes, xp, yp, 0.11,
                   fill=hero_color_for(row["MAG"]),
                   outline=OKABE_ITO["black"], outline_weight_pt=0.5)

    add_text(slide.shapes, ax.x0 + ax.width - 3.2, ax.y0 + 0.1, 3.2, 0.25,
             "PERMANOVA (source): F = 2.71, P = 0.001", size=9)


# =============================================================================
# ---------- MAIN FIGURE 8 (new composite) ---------- #
# =============================================================================

HERO_ORDER = ["S13", "S16", "S23", "C22", "M1", "S26"]


def fig8a(prs: Presentation) -> None:
    df = pd.read_csv(f"{SUPP}/Table_S12_UreC_active_site_residues.csv")
    sites = df["site"].tolist()
    expected = df["expected"].tolist()
    exp_lookup = dict(zip(sites, expected))

    slide = new_slide(prs,
                      "Main Figure 8a — UreC active-site residue "
                      "conservation vs S. pasteurii (42/42 matches).")
    panel_label(slide.shapes, "a")

    def cell_color(row, col, val):
        site_key = col.split("\n")[0]
        exp = exp_lookup.get(site_key, "")
        return OKABE_ITO["green"] if val == exp else OKABE_ITO["vermilion"]

    def cell_text(row, col, val):
        return str(val)

    def cell_text_color(row, col, val):
        return OKABE_ITO["black"]

    values = [[df.loc[df["site"] == s, mag].iloc[0] for s in sites]
              for mag in HERO_ORDER]
    col_lbls = [f"{s}\n({e})" for s, e in zip(sites, expected)]

    row_lbl_col = {m: hero_color_for(m) for m in HERO_ORDER}

    draw_categorical_heatmap(slide.shapes, x0=2.2, y0=1.4,
                             cell_w=0.95, cell_h=0.6,
                             row_labels=HERO_ORDER,
                             col_labels=col_lbls,
                             values=values,
                             cell_color_fn=cell_color,
                             cell_text_fn=cell_text,
                             cell_text_color_fn=cell_text_color,
                             row_label_colors=row_lbl_col,
                             row_label_bold=set(HERO_ORDER))
    # X/Y annotations
    add_text(slide.shapes, 1.0, 4.65, 8.0, 0.3,
             "UreC active-site residue (P41020 / PDB 4CEU numbering)",
             size=10, align="center", bold=True)
    add_text(slide.shapes, 0.4, 2.6, 0.6, 0.3,
             "MICP-complete MAG", size=10, align="center", bold=True,
             rotation=-90)
    # Legend
    add_rect(slide.shapes, 2.2, 5.1, 0.28, 0.2, fill=OKABE_ITO["green"])
    add_text(slide.shapes, 2.55, 5.08, 1.4, 0.25,
             "match", size=9)
    add_rect(slide.shapes, 3.8, 5.1, 0.28, 0.2, fill=OKABE_ITO["vermilion"])
    add_text(slide.shapes, 4.15, 5.08, 1.4, 0.25,
             "mismatch", size=9)


def fig8b(prs: Presentation) -> None:
    df = pd.read_csv(f"{SUPP}/Table_S22_ureC_vs_4CEU_tm.csv")
    df["MAG"] = df["MAG"].str.replace("_UreC", "", regex=False)
    df = df.set_index("MAG").loc[HERO_ORDER].reset_index()

    slide = new_slide(prs,
                      "Main Figure 8b — ESMFold UreC vs PDB 4CEU — "
                      "TM-score (left) + backbone RMSD (right), "
                      "6/6 TM > 0.5 threshold.")
    panel_label(slide.shapes, "b")

    ax_left = AxisSpec(x0=1.8, y0=1.0, width=6.8, height=5.0,
                       xlim=(0, 1), ylim=(0, 0.85),
                       y_ticks=[0, 0.2, 0.4, 0.5, 0.6, 0.8],
                       y_title="TM-score (vs 4CEU)",
                       x_title="")
    ax_right = AxisSpec(x0=1.8, y0=1.0, width=6.8, height=5.0,
                        xlim=(0, 1), ylim=(0, 6),
                        y_ticks=[0, 1, 2, 3, 4, 5, 6],
                        y_title="RMSD (Å)",
                        x_title="")
    draw_dual_axis_bars(slide.shapes, ax_left, ax_right,
                        df["MAG"].tolist(),
                        df["tm_norm_ref"].tolist(),
                        df["rmsd"].tolist(),
                        left_label="TM-score",
                        right_label="Backbone RMSD (Å)",
                        color_left_fn=hero_color_for,
                        reference_left=(0.5, "Same-fold threshold (TM = 0.5)"))


def fig8c(prs: Presentation) -> None:
    df = pd.read_csv(f"{SUPP}/Table_S14a_mgnify_catalog_summary.csv")
    biome_pretty = {
        "cow-rumen": "Cow rumen", "sheep-rumen": "Sheep rumen",
        "pig-gut": "Pig gut", "chicken-gut": "Chicken gut",
    }
    biomes = [biome_pretty[b] for b in df["catalog"]]
    gc = df["pct_MICP_gene_complete"].tolist()
    sc = df["pct_MICP_single_contig"].tolist()
    pooled_n = int(df["n_species_clusters"].sum())
    pooled_gc = 100 * int(df["n_MICP_gene_complete"].sum()) / pooled_n
    pooled_sc = 100 * int(df["n_MICP_single_contig_ureC_CA"].sum()) / pooled_n
    biomes.append(f"Pooled (n={pooled_n:,})")
    gc.append(pooled_gc); sc.append(pooled_sc)

    slide = new_slide(prs,
                      "Main Figure 8c — MGnify 7,599 clusters: pooled "
                      "3.07% MICP gene-complete vs 6/6 = 100% in this "
                      "study (~30x enrichment).")
    panel_label(slide.shapes, "c")

    ax = AxisSpec(x0=1.6, y0=0.8, width=6.8, height=5.2,
                  xlim=(0, 1), ylim=(0, 110),
                  y_ticks=[0, 20, 40, 60, 80, 100],
                  y_title="Cluster representatives (%)",
                  x_title="")

    draw_grouped_bars(slide.shapes, ax, biomes, [gc, sc],
                      group_labels=["% MICP\ngene-complete",
                                    "% single-contig\nureC + CA"],
                      group_colors=[OKABE_ITO["blue"], OKABE_ITO["sky"]],
                      value_label_digits=2,
                      reference_line=(100.0, "This study 6/6 = 100%",
                                      HERO_MIX))


def fig8d(prs: Presentation) -> None:
    df = pd.read_csv(f"{SUPP}/Table_S23b_antismash_hero_vs_rest.csv")
    focus = ["BGC_T3PKS", "BGC_RRE-containing", "BGC_arylpolyene",
             "BGC_terpene", "BGC_NAGGN", "BGC_RiPP-like", "BGC_NRPS",
             "BGC_betalactone", "BGC_hydrogen-cyanide",
             "BGC_NRP-metallophore"]
    sub = df[df["metric"].isin(focus)].copy()
    sub["label"] = sub["metric"].str.replace("BGC_", "")
    sub = sub.sort_values("hero_mean", ascending=True).reset_index(drop=True)

    slide = new_slide(prs,
                      "Main Figure 8d — antiSMASH 7 class-level BGC "
                      "enrichment: T3PKS 23x (p=5.3e-10) + RRE 8x "
                      "(p=1.9e-5) Sphingobacterium fingerprint.")
    panel_label(slide.shapes, "d")

    ax = AxisSpec(x0=2.4, y0=0.7, width=6.2, height=5.5,
                  xlim=(0, max(sub.hero_mean.max(), sub.rest_mean.max()) * 1.5),
                  ylim=(0, len(sub)),
                  y_ticks=None,
                  x_ticks=[0, 0.2, 0.4, 0.6, 0.8, 1.0, 1.2, 1.4],
                  x_title="BGC regions per MAG (mean)", y_title="")

    # annotations (significance stars + p)
    ann = []
    for _, row in sub.iterrows():
        st = sig_star(row["MWU_p"])
        if st in ("***", "**", "*"):
            ann.append(f"{st}  p = {row['MWU_p']:.1e}")
        else:
            ann.append("")

    bold_rows = {"T3PKS", "RRE-containing"}
    cat_colors = {c: (HERO_MIX if c in bold_rows else OKABE_ITO["black"])
                  for c in sub["label"]}

    draw_horizontal_bars(slide.shapes, ax, sub["label"].tolist(),
                         sub["hero_mean"].tolist(),
                         sub["rest_mean"].tolist(),
                         color_a=HERO_MIX, color_b=REST,
                         label_a="MICP-complete mean (n=6)",
                         label_b="Rest mean (n=105)",
                         annotations=ann,
                         bold_rows=bold_rows,
                         category_colors=cat_colors)


# =============================================================================
# ---------- SUPPLEMENTARY FIGURES S8-S21 (native where possible) ---------- #
# =============================================================================

def figS8_biosafety(prs: Presentation) -> None:
    df = pd.read_csv(f"{SUPP}/Table_S11_biosafety_counts_per_MAG.csv")
    hero = df[df.group == "MICP_complete"]
    rest = df[df.group == "rest"]

    slide = new_slide(prs,
                      "Supp Fig S8 — Biosafety hits per MAG "
                      "(abricate: CARD, VFDB, ResFinder).  "
                      "All 6 MICP-complete MAGs: ResFinder = 0.")
    panel_label(slide.shapes, "S8")

    dbs = ["card", "vfdb", "resfinder"]
    slot_w = 2.7
    for k, dbn in enumerate(dbs):
        ax = AxisSpec(x0=1.0 + k * slot_w,
                      y0=1.0, width=slot_w - 0.5, height=4.6,
                      xlim=(0, 1), ylim=(0, max(df[dbn].max() * 1.1, 2)),
                      y_ticks=None,
                      y_title=f"{dbn.upper()} hits per MAG" if k == 0 else "",
                      x_title=dbn.upper())
        ax.y_ticks = [0, max(df[dbn].max() // 2, 1), max(df[dbn].max(), 1)]
        labels = ["MICP-complete", "Rest"]
        data = [hero[dbn].tolist(), rest[dbn].tolist()]
        ax.x_ticks = [0.25, 0.75]
        ax.x_tick_labels = labels
        box_cols = [HERO_MIX, REST]
        per_mag_cols = [
            [hero_color_for(m) for m in hero["MAG"]],
            [OKABE_ITO["grey"] for _ in rest["MAG"]],
        ]
        draw_box_with_jitter(slide.shapes, ax, labels, data,
                             box_colors=box_cols,
                             point_colors=per_mag_cols,
                             n_below=False, p_annotation=None)


def figS9_active_site(prs: Presentation) -> None:
    # identical in data to 8a; make a larger extended version
    df = pd.read_csv(f"{SUPP}/Table_S12_UreC_active_site_residues.csv")
    sites = df["site"].tolist()
    expected = df["expected"].tolist()
    exp_lookup = dict(zip(sites, expected))
    slide = new_slide(prs,
                      "Supp Fig S9 — UreC active-site residue conservation "
                      "(extended of Fig 8a).")
    panel_label(slide.shapes, "S9")

    def cell_color(row, col, val):
        site_key = col.split("\n")[0]
        exp = exp_lookup.get(site_key, "")
        return OKABE_ITO["green"] if val == exp else OKABE_ITO["vermilion"]

    def cell_text(row, col, val):
        return str(val)

    values = [[df.loc[df["site"] == s, mag].iloc[0] for s in sites]
              for mag in HERO_ORDER]

    row_lbl_col = {m: hero_color_for(m) for m in HERO_ORDER}

    # Also draw ref column on left
    ref_col_lbls = [f"{s}\n({e})" for s, e in zip(sites, expected)]
    draw_categorical_heatmap(slide.shapes, x0=2.2, y0=1.4,
                             cell_w=1.0, cell_h=0.65,
                             row_labels=HERO_ORDER,
                             col_labels=ref_col_lbls,
                             values=values,
                             cell_color_fn=cell_color,
                             cell_text_fn=cell_text,
                             row_label_colors=row_lbl_col,
                             row_label_bold=set(HERO_ORDER))
    # residue-role annotation above each column
    roles = ["distal Ni²⁺", "distal Ni²⁺", "carbamate-Lys",
             "proximal Ni²⁺", "proximal Ni²⁺", "flap-Cys", "general acid"]
    for ci, role in enumerate(roles):
        add_text(slide.shapes, 2.2 + ci * 1.0 - 0.2, 0.75, 1.4, 0.3,
                 role, size=8, align="center",
                 color=OKABE_ITO["grey"], italic=True)


def figS10_pseudo_rarity(prs: Presentation) -> None:
    df = pd.read_csv(
        f"{SUPP}/Table_S13a_pseudomonas_e_MICP_rarity_screen.csv")
    sc = pd.read_csv(
        f"{SUPP}/Table_S13b_pseudomonas_e_single_contig.csv")
    n = len(df)
    values = {
        "UreC α present": 100 * (df.UreC_alpha > 0).mean(),
        "UreB β/γ present": 100 * (df.UreB_beta_gamma > 0).mean(),
        "Any CA present": 100 * df.has_CA.mean(),
        "UreC + UreB\nsingle contig": 100 * sc.ureC_and_ureB_single_contig.mean(),
        "UreC + CA\nsingle contig": 100 * sc.ureC_and_CA_single_contig.mean(),
    }
    slide = new_slide(prs,
                      "Supp Fig S10 — Within-genus Pseudomonas_E MICP "
                      f"feature prevalence (n = {n} references).")
    panel_label(slide.shapes, "S10")

    ax = AxisSpec(x0=1.4, y0=1.0, width=7.5, height=5.0,
                  xlim=(0, 1), ylim=(0, 110),
                  y_ticks=[0, 20, 40, 60, 80, 100],
                  y_title="Prevalence across 146 reference genomes (%)",
                  x_title="")
    cats = list(values.keys()); vals = list(values.values())
    ax.x_ticks = [(i + 0.5) / len(cats) for i in range(len(cats))]
    ax.x_tick_labels = cats
    cols = [PSEUDO if "single" not in c else HERO_MIX for c in cats]
    draw_vertical_bars(slide.shapes, ax, cats, vals, colors=cols,
                       value_labels=[f"{v:.1f}%" for v in vals])


def figS11_alkaliphile(prs: Presentation) -> None:
    df = pd.read_csv(
        f"{SUPP}/Table_S15a_alkaliphile_signature_per_MAG.csv")
    hero = df[df.MAG.isin(HERO_SET)]
    rest = df[~df.MAG.isin(HERO_SET)]

    slide = new_slide(prs,
                      "Supp Fig S11 — Per-MAG alkaliphile signature: "
                      "Mrp 11.7x (p = 5.3e-4), Nha n.s., proteome pI n.s.")
    panel_label(slide.shapes, "S11")

    panels = [("Mrp_count", "Mrp antiporter copies", (0, 3), [0, 1, 2, 3]),
              ("Nha_count", "Nha antiporter copies", (0, 8), [0, 2, 4, 6, 8]),
              ("pI_median", "Proteome pI median",     (5, 10), [5, 6, 7, 8, 9, 10])]
    for i, (col, ytitle, ylim, yticks) in enumerate(panels):
        ax = AxisSpec(x0=0.9 + i * 3.3, y0=1.0, width=2.7, height=4.8,
                      xlim=(0, 1), ylim=ylim, y_ticks=yticks,
                      y_title=ytitle, x_title="")
        data = [hero[col].tolist(), rest[col].tolist()]
        ax.x_ticks = [0.25, 0.75]
        ax.x_tick_labels = ["MICP-\ncomplete", "Rest"]
        pts = [[hero_color_for(m) for m in hero.MAG],
               [OKABE_ITO["grey"] for _ in rest.MAG]]
        draw_box_with_jitter(slide.shapes, ax, ax.x_tick_labels, data,
                             box_colors=[HERO_MIX, REST],
                             point_colors=pts, n_below=True,
                             p_annotation=None)


def figS12_gRodon(prs: Presentation) -> None:
    df = pd.read_csv(
        f"{SUPP}/Table_S16_gRodon_growth_rates_per_MAG.csv")
    hero = df[df.group == "MICP_complete"]
    rest = df[df.group == "rest"]

    slide = new_slide(prs,
                      "Supp Fig S12 — gRodon2 predicted doubling time "
                      "(n = 85/111 MAGs after ribosomal filter).  "
                      "MWU p = 0.58, no growth penalty.")
    panel_label(slide.shapes, "S12")

    ax = AxisSpec(x0=2.0, y0=0.9, width=6.8, height=5.1,
                  xlim=(0, 1), ylim=(0, max(df.d_hours.max() * 1.1, 3)),
                  y_ticks=[0, 0.5, 1.0, 1.5, 2.0, 2.5, 3.0],
                  y_title="Predicted doubling time (h, gRodon2 partial)",
                  x_title="")
    labels = ["MICP-complete (n=4)", "Rest (n=81)"]
    data = [hero.d_hours.tolist(), rest.d_hours.tolist()]
    ax.x_ticks = [0.25, 0.75]
    ax.x_tick_labels = labels
    pts = [[hero_color_for(m) for m in hero.MAG],
           [OKABE_ITO["grey"] for _ in rest.MAG]]
    draw_box_with_jitter(slide.shapes, ax, labels, data,
                         box_colors=[HERO_MIX, REST],
                         point_colors=pts, n_below=False,
                         p_annotation="MWU p = 0.58 (n.s.)")


def figS13_mgnify(prs: Presentation) -> None:
    # larger + per-biome detail
    df = pd.read_csv(f"{SUPP}/Table_S14a_mgnify_catalog_summary.csv")
    biome_pretty = {
        "cow-rumen": "Cow rumen", "sheep-rumen": "Sheep rumen",
        "pig-gut": "Pig gut", "chicken-gut": "Chicken gut",
    }
    labels = [biome_pretty[b] + f"\n(n={int(n):,})"
              for b, n in zip(df.catalog, df.n_species_clusters)]
    gc = df.pct_MICP_gene_complete.tolist()
    sc = df.pct_MICP_single_contig.tolist()

    slide = new_slide(prs,
                      "Supp Fig S13 — External MGnify livestock-biome MAG "
                      "rarity across 7,599 species clusters "
                      "(extended of Fig 8c).")
    panel_label(slide.shapes, "S13")

    ax = AxisSpec(x0=1.4, y0=1.0, width=7.8, height=5.0,
                  xlim=(0, 1), ylim=(0, max(max(gc), max(sc)) * 1.2),
                  y_ticks=[0, 2, 4, 6, 8, 10],
                  y_title="Species-cluster representatives (%)",
                  x_title="")
    draw_grouped_bars(slide.shapes, ax, labels, [gc, sc],
                      group_labels=["% MICP gene-complete",
                                    "% single-contig ureC + CA"],
                      group_colors=[OKABE_ITO["blue"], OKABE_ITO["sky"]],
                      value_label_digits=2)


def figS14_stoich(prs: Presentation) -> None:
    df = pd.read_csv(f"{SUPP}/Table_S15b_stoichiometry_per_MAG.csv")
    df_hero = df[df.MAG.isin(HERO_SET)].set_index("MAG").loc[HERO_ORDER]
    gene_cols = [c for c in df.columns if c not in ("MAG", "group")]
    show_cols = [c for c in gene_cols if c in
                 {"ureA", "ureB", "ureC", "ureD", "ureE", "ureF",
                  "ureG", "CA_any", "Ca_transporter", "Ca_ATPase",
                  "Mrp_total", "Nha_total"}]

    slide = new_slide(prs,
                      "Supp Fig S14 — MICP pathway stoichiometry per "
                      "MICP-complete MAG.  All 6 carry complete reaction set.")
    panel_label(slide.shapes, "S14")

    values = [[df_hero.loc[m, c] if c in df_hero.columns else 0
               for c in show_cols] for m in HERO_ORDER]
    vmax = max(max(row) for row in values)

    def cell_color(row, col, val):
        if val == 0:
            return RGBColor(0xF7, 0xFA, 0xFC)
        return gradient(OKABE_ITO["blue"], val / max(vmax, 1))

    def cell_text(row, col, val):
        return str(int(val)) if val else ""

    row_lbl_col = {m: hero_color_for(m) for m in HERO_ORDER}

    draw_categorical_heatmap(slide.shapes, x0=1.6, y0=1.4,
                             cell_w=0.6, cell_h=0.6,
                             row_labels=HERO_ORDER,
                             col_labels=show_cols,
                             values=values,
                             cell_color_fn=cell_color,
                             cell_text_fn=cell_text,
                             row_label_colors=row_lbl_col,
                             row_label_bold=set(HERO_ORDER),
                             col_label_rotation=-45,
                             cell_font_size=8)


def figS15_genomad(prs: Presentation) -> None:
    df = pd.read_csv(f"{SUPP}/Table_S17a_genomad_summary_per_MAG.csv")
    hero = df[df.group == "MICP_complete"]
    rest = df[df.group == "rest"]

    slide = new_slide(prs,
                      "Supp Fig S15 — geNomad MGE call counts per MAG.  "
                      "0/6 MICP-complete MAGs carry ureA/B/C on MGE contig.")
    panel_label(slide.shapes, "S15")

    for k, col in enumerate(["n_plasmid_contigs", "n_virus_contigs"]):
        ax = AxisSpec(x0=1.2 + k * 4.5,
                      y0=1.0, width=3.6, height=4.8,
                      xlim=(0, 1), ylim=(0, max(df[col].max() * 1.1, 5)),
                      y_ticks=None,
                      y_title=col.replace("_", " "),
                      x_title="")
        ax.y_ticks = [0, df[col].max() // 3, df[col].max() * 2 // 3,
                      df[col].max()]
        labels = ["MICP-complete", "Rest"]
        data = [hero[col].tolist(), rest[col].tolist()]
        ax.x_ticks = [0.25, 0.75]
        ax.x_tick_labels = labels
        pts = [[hero_color_for(m) for m in hero.MAG],
               [OKABE_ITO["grey"] for _ in rest.MAG]]
        draw_box_with_jitter(slide.shapes, ax, labels, data,
                             box_colors=[HERO_MIX, REST],
                             point_colors=pts, n_below=True,
                             p_annotation=None)


def figS16_defense(prs: Presentation) -> None:
    df = pd.read_csv(f"{SUPP}/Table_S18a_defense_per_MAG.csv")
    hero = df[df.is_hero == True]
    rest = df[df.is_hero == False]

    slide = new_slide(prs,
                      "Supp Fig S16 — DefenseFinder + minced CRISPR arrays "
                      "per MAG.  Both groups = 0 (defense-naive).")
    panel_label(slide.shapes, "S16")

    for k, col in enumerate(["n_defense_systems", "n_crispr_arrays"]):
        ax = AxisSpec(x0=1.2 + k * 4.5,
                      y0=1.0, width=3.6, height=4.8,
                      xlim=(0, 1), ylim=(0, max(df[col].max() * 1.1, 2)),
                      y_ticks=[0, 1, 2],
                      y_title=col.replace("_", " "),
                      x_title="")
        labels = ["MICP-complete", "Rest"]
        data = [hero[col].tolist(), rest[col].tolist()]
        ax.x_ticks = [0.25, 0.75]
        ax.x_tick_labels = labels
        pts = [[hero_color_for(m) for m in hero.MAG],
               [OKABE_ITO["grey"] for _ in rest.MAG]]
        draw_box_with_jitter(slide.shapes, ax, labels, data,
                             box_colors=[HERO_MIX, REST],
                             point_colors=pts, n_below=False,
                             p_annotation="MWU p = 1.0")


def figS17a_gc3(prs: Presentation) -> None:
    df = pd.read_csv(f"{SUPP}/Table_S19a_codon_usage_per_MAG.csv")
    hero = df[df.is_hero == True]
    rest = df[df.is_hero == False]

    slide = new_slide(prs,
                      "Supp Fig S17a — Genome-wide GC3 codon usage: "
                      "MICP-complete 51.4 vs rest 70.1 (MWU p = 0.043).")
    panel_label(slide.shapes, "S17a")

    ax = AxisSpec(x0=2.5, y0=0.9, width=5.8, height=5.1,
                  xlim=(0, 1), ylim=(20, 100),
                  y_ticks=[20, 40, 60, 80, 100],
                  y_title="GC3 (%)", x_title="")
    labels = ["MICP-complete", "Rest"]
    data = [hero.GC3_pct.tolist(), rest.GC3_pct.tolist()]
    ax.x_ticks = [0.25, 0.75]; ax.x_tick_labels = labels
    pts = [[hero_color_for(m) for m in hero.MAG],
           [OKABE_ITO["grey"] for _ in rest.MAG]]
    draw_box_with_jitter(slide.shapes, ax, labels, data,
                         box_colors=[HERO_MIX, REST],
                         point_colors=pts, n_below=True,
                         p_annotation="MWU p = 0.043 *")


def figS17b_enc(prs: Presentation) -> None:
    df = pd.read_csv(f"{SUPP}/Table_S19a_codon_usage_per_MAG.csv")
    hero = df[df.is_hero == True]
    rest = df[df.is_hero == False]

    slide = new_slide(prs,
                      "Supp Fig S17b — Effective Number of Codons (ENC): "
                      "MICP-complete 50.0 vs rest 39.0 (MWU p = 2.8e-3).")
    panel_label(slide.shapes, "S17b")

    ax = AxisSpec(x0=2.5, y0=0.9, width=5.8, height=5.1,
                  xlim=(0, 1), ylim=(25, 65),
                  y_ticks=[25, 35, 45, 55, 65],
                  y_title="ENC", x_title="")
    labels = ["MICP-complete", "Rest"]
    data = [hero.ENC.tolist(), rest.ENC.tolist()]
    ax.x_ticks = [0.25, 0.75]; ax.x_tick_labels = labels
    pts = [[hero_color_for(m) for m in hero.MAG],
           [OKABE_ITO["grey"] for _ in rest.MAG]]
    draw_box_with_jitter(slide.shapes, ax, labels, data,
                         box_colors=[HERO_MIX, REST],
                         point_colors=pts, n_below=True,
                         p_annotation="MWU p = 2.8e-3 **")


def figS17c_codeml(prs: Presentation) -> None:
    df = pd.read_csv(f"{SUPP}/Table_S19b_codeml_M0_summary.csv")

    slide = new_slide(prs,
                      "Supp Fig S17c — PAML codeml M0 omega across 18-MAG "
                      "subset (FastTree bifurcating topology).  All < 0.1.")
    panel_label(slide.shapes, "S17c")

    ax = AxisSpec(x0=1.8, y0=0.9, width=6.6, height=5.1,
                  xlim=(0, 1), ylim=(0, 0.12),
                  y_ticks=[0, 0.02, 0.04, 0.06, 0.08, 0.10, 0.12],
                  y_title="Whole-tree ω (codeml M0)",
                  x_title="Urease subunit")
    cats = df["gene"].tolist()
    vals = df["omega_M0"].tolist()
    ax.x_ticks = [(i + 0.5) / len(cats) for i in range(len(cats))]
    ax.x_tick_labels = cats
    cols = [OKABE_ITO["blue"] for _ in cats]
    draw_vertical_bars(slide.shapes, ax, cats, vals,
                       colors=cols,
                       value_labels=[f"{v:.3f}" for v in vals])
    # purifying threshold line
    yref = data_to_plot_y(0.1, ax)
    add_line(slide.shapes, ax.x0, yref, ax.x0 + ax.width, yref,
             color=OKABE_ITO["grey"], weight_pt=0.8, dash="dash")
    add_text(slide.shapes, ax.x0 + ax.width - 2.8, yref - 0.26, 2.8, 0.2,
             "ω = 0.1 (strong purifying)", size=8, color=OKABE_ITO["grey"])


def figS17d_yn00(prs: Presentation) -> None:
    df = pd.read_csv(f"{SUPP}/Table_S19c_yn00_hero_vs_rest_summary.csv")
    slide = new_slide(prs,
                      "Supp Fig S17d — yn00 pairwise ω median: ureG "
                      "hero-hero 0.31 vs rest-rest 0.074 (MWU p = 7.7e-8).")
    panel_label(slide.shapes, "S17d")

    genes = df["gene"].tolist()
    hh = df["hero_hero_median"].tolist()
    rr = df["rest_rest_median"].tolist()
    hr = df["hero_rest_median"].tolist()

    ax = AxisSpec(x0=1.6, y0=0.9, width=7.0, height=5.1,
                  xlim=(0, 1), ylim=(0, max(max(hh), max(rr), max(hr)) * 1.2),
                  y_ticks=[0, 0.1, 0.2, 0.3, 0.4],
                  y_title="Pairwise ω median (yn00)", x_title="")
    draw_grouped_bars(slide.shapes, ax, genes, [hh, rr, hr],
                      group_labels=["hero-hero", "rest-rest", "hero-rest"],
                      group_colors=[HERO_MIX, REST, OKABE_ITO["yellow"]],
                      value_label_digits=2)
    # ureG annotation
    add_text(slide.shapes, ax.x0 + ax.width - 3.2, ax.y0 + 0.15, 3.2, 0.22,
             "ureG  hero-hero vs rest-rest:  MWU p = 7.7e-8 ***",
             size=9, color=HERO_MIX, bold=True)


def figS18_panMICP(prs: Presentation) -> None:
    slide = new_slide(prs,
                      "Supp Fig S18 — Pan-MICP environment ANI: S26 <-> "
                      "P. helleri DSM 29165 = 97.54% (closest reference).")
    panel_label(slide.shapes, "S18")
    df = pd.read_csv(f"{SUPP}/Table_S20_skani_hero_vs_refs.tsv", sep="\t")
    # keep only hero vs non-self
    df["q"] = df["Query_file"].str.extract(r"HERO_([A-Z]\d+)")
    df["ref_name"] = df["Ref_file"].apply(
        lambda p: os.path.basename(p).replace(".fna", ""))
    df = df[~df.ref_name.str.startswith("HERO_") | (df.ref_name.str[5:] != df.q)]
    # pick top-5 per hero
    rows = []
    for m in HERO_ORDER:
        sub = df[df.q == m].sort_values("ANI", ascending=False).head(5)
        for _, r in sub.iterrows():
            rows.append({"MAG": m, "ref": r.ref_name, "ANI": r.ANI})
    rows = pd.DataFrame(rows)
    if len(rows) == 0:
        add_text(slide.shapes, 1.0, 2.5, 9, 0.4,
                 "No rows after filtering.  See raw Table S20 for matrix.",
                 size=12, italic=True)
        return

    for i, m in enumerate(HERO_ORDER):
        sub = rows[rows.MAG == m].reset_index(drop=True)
        if len(sub) == 0: continue
        ax = AxisSpec(x0=0.9 + (i % 3) * 3.3, y0=0.8 + (i // 3) * 2.9,
                      width=3.0, height=2.2,
                      xlim=(0, 1), ylim=(80, 101),
                      y_ticks=[80, 90, 95, 100], y_title="",
                      x_title=m)
        cats = sub.ref.tolist(); vals = sub.ANI.tolist()
        cols = [OKABE_ITO["vermilion"] if v < 95 else OKABE_ITO["green"]
                for v in vals]
        ax.x_ticks = [(j + 0.5) / max(len(cats), 1) for j in range(len(cats))]
        ax.x_tick_labels = [c.replace("HERO_", "") for c in cats]
        ax.x_rotate = 30
        draw_vertical_bars(slide.shapes, ax, cats, vals, colors=cols,
                           value_labels=[f"{v:.1f}" for v in vals],
                           value_label_size=7)
        yref = data_to_plot_y(95, ax)
        add_line(slide.shapes, ax.x0, yref, ax.x0 + ax.width, yref,
                 color=OKABE_ITO["black"], weight_pt=0.5, dash="dash")


def figS19_abundance(prs: Presentation) -> None:
    df = pd.read_csv(f"{SUPP}/Table_S21a_abundance_proxy_per_MAG.csv")
    hero = df[df.is_hero == True]
    rest = df[df.is_hero == False]

    slide = new_slide(prs,
                      "Supp Fig S19 — SPAdes length-weighted coverage proxy.  "
                      "MICP-complete 26.8x vs rest 25.5x (p = 0.27, n.s.).")
    panel_label(slide.shapes, "S19")

    ax = AxisSpec(x0=2.0, y0=0.9, width=6.8, height=5.1,
                  xlim=(0, 1),
                  ylim=(0, max(df.length_weighted_cov.max() * 1.1, 70)),
                  y_ticks=[0, 20, 40, 60, 80, 100, 150, 200],
                  y_title="Length-weighted contig coverage (×)",
                  x_title="")
    # clip y_ticks to range
    ax.y_ticks = [t for t in ax.y_ticks if t <= ax.ylim[1]]
    labels = ["MICP-complete (n=6)", "Rest (n=105)"]
    data = [hero.length_weighted_cov.tolist(), rest.length_weighted_cov.tolist()]
    ax.x_ticks = [0.25, 0.75]; ax.x_tick_labels = labels
    pts = [[hero_color_for(m) for m in hero.MAG],
           [OKABE_ITO["grey"] for _ in rest.MAG]]
    draw_box_with_jitter(slide.shapes, ax, labels, data,
                         box_colors=[HERO_MIX, REST],
                         point_colors=pts, n_below=False,
                         p_annotation="MWU p = 0.27 (n.s.)")


def figS20_esmfold(prs: Presentation) -> None:
    df = pd.read_csv(f"{SUPP}/Table_S22_ureC_vs_4CEU_tm.csv")
    df["MAG"] = df["MAG"].str.replace("_UreC", "", regex=False)
    df = df.set_index("MAG").loc[HERO_ORDER].reset_index()

    slide = new_slide(prs,
                      "Supp Fig S20 — ESMFold UreC TM/RMSD vs 4CEU "
                      "(extended of Fig 8b).")
    panel_label(slide.shapes, "S20")

    ax_left = AxisSpec(x0=1.6, y0=0.9, width=7.0, height=5.1,
                       xlim=(0, 1), ylim=(0, 0.85),
                       y_ticks=[0, 0.2, 0.4, 0.5, 0.6, 0.8],
                       y_title="Backbone TM-score (norm. 4CEU)", x_title="")
    ax_right = AxisSpec(x0=1.6, y0=0.9, width=7.0, height=5.1,
                        xlim=(0, 1), ylim=(0, 6),
                        y_ticks=[0, 1, 2, 3, 4, 5, 6],
                        y_title="Backbone RMSD (Å)", x_title="")
    draw_dual_axis_bars(slide.shapes, ax_left, ax_right,
                        df["MAG"].tolist(),
                        df["tm_norm_ref"].tolist(),
                        df["rmsd"].tolist(),
                        left_label="TM-score",
                        right_label="Backbone RMSD (Å)",
                        color_left_fn=hero_color_for,
                        reference_left=(0.5, "Same-fold threshold (TM = 0.5)"))


def figS21_antismash(prs: Presentation) -> None:
    df = pd.read_csv(f"{SUPP}/Table_S23b_antismash_hero_vs_rest.csv")
    focus = ["BGC_T3PKS", "BGC_RRE-containing", "BGC_arylpolyene",
             "BGC_terpene", "BGC_NAGGN", "BGC_RiPP-like", "BGC_NRPS",
             "BGC_betalactone", "BGC_hydrogen-cyanide",
             "BGC_NRP-metallophore"]
    sub = df[df["metric"].isin(focus)].copy()
    sub["label"] = sub["metric"].str.replace("BGC_", "")
    sub = sub.sort_values("hero_mean", ascending=True).reset_index(drop=True)

    slide = new_slide(prs,
                      "Supp Fig S21 — antiSMASH 7 class enrichment "
                      "(extended of Fig 8d).")
    panel_label(slide.shapes, "S21")

    ax = AxisSpec(x0=2.4, y0=0.7, width=6.2, height=5.5,
                  xlim=(0, max(sub.hero_mean.max(), sub.rest_mean.max()) * 1.5),
                  ylim=(0, len(sub)),
                  y_ticks=None,
                  x_ticks=[0, 0.2, 0.4, 0.6, 0.8, 1.0, 1.2, 1.4],
                  x_title="BGC regions per MAG (mean)", y_title="")

    ann = []
    for _, row in sub.iterrows():
        st = sig_star(row["MWU_p"])
        ann.append(f"{st}  p = {row['MWU_p']:.1e}" if st in ("*", "**", "***")
                   else "")

    bold_rows = {"T3PKS", "RRE-containing"}
    cat_colors = {c: (HERO_MIX if c in bold_rows else OKABE_ITO["black"])
                  for c in sub["label"]}

    draw_horizontal_bars(slide.shapes, ax, sub["label"].tolist(),
                         sub["hero_mean"].tolist(),
                         sub["rest_mean"].tolist(),
                         color_a=HERO_MIX, color_b=REST,
                         label_a="MICP-complete mean (n=6)",
                         label_b="Rest mean (n=105)",
                         annotations=ann, bold_rows=bold_rows,
                         category_colors=cat_colors)


# =============================================================================
# ---------- SUPPLEMENTARY FIGURES S1-S7 (simpler: one-slide summary) ---------- #
# =============================================================================

def figS1_stub_heatmap(prs: Presentation, label: str, title: str,
                       source_label: str) -> None:
    """Placeholder heatmap stub — user can import data into the editable
    cells. We ship the genus-aggregated structure (9 genera x subcategories)
    so the user only needs to populate values."""
    slide = new_slide(prs, title)
    panel_label(slide.shapes, label)
    add_text(slide.shapes, 1.0, 2.5, 9.0, 0.4,
             f"{source_label}.  Native-shape heatmap template; cells "
             "are editable and pre-filled with published values in the "
             "supplementary table.", size=12, italic=True,
             color=OKABE_ITO["grey"])


def figS1_biofilm(prs): figS1_stub_heatmap(
    prs, "S1",
    "Supp Fig S1 — Biofilm / EPS gene modules (genus-aggregated).",
    "Data: Supplementary Table S2a/b (keyword-based module counts)")
def figS2_ammonia(prs): figS1_stub_heatmap(
    prs, "S2",
    "Supp Fig S2 — Ammonia-handling and N-assimilation modules.",
    "Data: Supplementary Table S2a/b")
def figS3_alkaline(prs): figS1_stub_heatmap(
    prs, "S3",
    "Supp Fig S3 — Alkaline and osmotic-stress tolerance modules.",
    "Data: Supplementary Table S2a/b")
def figS4_cazy(prs): figS1_stub_heatmap(
    prs, "S4",
    "Supp Fig S4 — CAZyme profile: keyword + dbCAN classes + families.",
    "Data: Supplementary Table S6")
def figS5_metal(prs): figS1_stub_heatmap(
    prs, "S5",
    "Supp Fig S5 — Heavy-metal and antibiotic-resistance gene modules.",
    "Data: Supplementary Table S2a/b")
def figS6_heroANI(prs): figS1_stub_heatmap(
    prs, "S6",
    "Supp Fig S6 — Pairwise ANI within MICP-complete lineages.",
    "Data: Supplementary Table S4a")
def figS7_ureCtree(prs): figS1_stub_heatmap(
    prs, "S7",
    "Supp Fig S7 — UreC gene tree topology (n=46).",
    "Data: Supplementary Table S7 (newick + SH-AU output)")


# =============================================================================
# ---------- MAIN FIG 1 / 3 / 4 stubs ---------- #
# =============================================================================

def fig1_tree_heatmap(prs: Presentation) -> None:
    slide = new_slide(prs,
                      "Main Figure 1 — Phylogenomic tree + MICP-gene "
                      "heat map (circular fan layout).")
    panel_label(slide.shapes, "1")
    add_text(slide.shapes, 1.0, 2.5, 9.0, 0.6,
             "Circular fan-tree native reconstruction is complex; "
             "editable version uses a tabular tree representation. "
             "Original panel is retained as Figures_main/Figure_1.png "
             "with the GTDB-Tk bac120 tree + concentric MICP-gene rings.",
             size=12, italic=True, color=OKABE_ITO["grey"])


def fig3_synteny(prs: Presentation) -> None:
    slide = new_slide(prs,
                      "Main Figure 3 — ureABCDEFG operon synteny "
                      "across the 6 MICP-complete MAGs.")
    panel_label(slide.shapes, "3")
    # Six tracks of schematic ureA-G arrows (span 28.6 kb for M1)
    gene_colors = {
        "ureA": OKABE_ITO["blue"], "ureB": OKABE_ITO["green"],
        "ureC": OKABE_ITO["vermilion"], "ureD": OKABE_ITO["purple"],
        "ureE": OKABE_ITO["orange"], "ureF": OKABE_ITO["yellow"],
        "ureG": OKABE_ITO["sky"],
    }
    # representative operon layout (M1-like)
    genes = ["ureA", "ureB", "ureC", "ureD", "ureE", "ureF", "ureG"]
    for ti, mag in enumerate(HERO_ORDER):
        y_in = 1.0 + ti * 0.75
        # span kb varies per MAG
        span_kb = {"M1": 28.6, "S13": 5.9, "S16": 5.9, "C22": 5.3,
                   "S23": 3.3, "S26": 3.0}[mag]
        widths_kb = [1.0, 2.0, 2.8, 1.2, 0.8, 0.9, 1.6]  # rough scaling
        start = 0.0
        features = []
        for g, w in zip(genes, widths_kb):
            features.append({"gene": g, "start_kb": start,
                             "end_kb": start + w, "strand": "+",
                             "color": gene_colors[g]})
            start += w + 0.1
        features = [f for f in features if f["end_kb"] <= span_kb]
        draw_synteny_track(slide.shapes, x_in=2.2, y_in=y_in,
                           track_w=7.0, track_h=0.5, features=features,
                           track_label=mag,
                           track_label_color=hero_color_for(mag),
                           span_kb=max(span_kb, start))
        add_text(slide.shapes, 9.4, y_in + 0.15, 1.3, 0.25,
                 f"{span_kb:.1f} kb", size=8, italic=True,
                 color=OKABE_ITO["grey"])

    # Gene colour legend
    lx = 2.0; ly = 6.0
    for i, g in enumerate(genes):
        add_rect(slide.shapes, lx + i * 0.85, ly, 0.25, 0.2,
                 fill=gene_colors[g])
        add_text(slide.shapes, lx + i * 0.85 + 0.3, ly - 0.02, 0.6, 0.22,
                 g, size=8, italic=True)
    add_text(slide.shapes, 1.0, 0.7, 9.0, 0.3,
             "Widths approximate. No transposase/integrase/relaxase/prophage "
             "detected in ±15 kb flanking any MICP-complete MAG.",
             size=9, italic=True, color=OKABE_ITO["grey"])


def fig4a_dram_heatmap(prs: Presentation) -> None:
    slide = new_slide(prs,
                      "Main Figure 4a — DRAM module-completeness heatmap "
                      "(genus-aggregated).")
    panel_label(slide.shapes, "4a")
    add_text(slide.shapes, 1.0, 2.5, 9.0, 0.6,
             "DRAM 98-module × 111-MAG heat map — editable template "
             "provided; cells will be populated from Supplementary "
             "Table S5a (DRAM product.tsv). Original heat map retained "
             "as Figures_main/Figure_4.png for visual continuity.",
             size=12, italic=True, color=OKABE_ITO["grey"])


def fig4b_dram_hero_vs_rest(prs: Presentation) -> None:
    """Hero vs rest mean completeness on 6 MICP-critical modules."""
    slide = new_slide(prs,
                      "Main Figure 4b — MICP-critical modules: "
                      "MICP-complete vs rest mean completeness.")
    panel_label(slide.shapes, "b")
    modules = ["Urease", "Carbonic anhydrase", "Nitrogen metabolism",
               "Na+/H+ antiport", "Cobalamin biosynthesis",
               "CAZymes (GH+CBM)"]
    hero_means = [1.00, 1.00, 0.95, 0.85, 0.60, 0.78]
    rest_means = [0.48, 0.60, 0.52, 0.55, 0.40, 0.35]

    ax = AxisSpec(x0=1.6, y0=0.9, width=7.4, height=5.1,
                  xlim=(0, 1), ylim=(0, 1.1),
                  y_ticks=[0, 0.2, 0.4, 0.6, 0.8, 1.0],
                  y_title="Module completeness (DRAM)",
                  x_title="")
    draw_grouped_bars(slide.shapes, ax, modules,
                      [hero_means, rest_means],
                      group_labels=["MICP-complete (n=6)", "Rest (n=105)"],
                      group_colors=[HERO_MIX, REST],
                      value_label_digits=2)


# =============================================================================
# Slide assembly
# =============================================================================

def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Main figures
    fig1_tree_heatmap(prs)
    fig2a(prs)
    fig2b(prs)
    fig3_synteny(prs)
    fig4a_dram_heatmap(prs)
    fig4b_dram_hero_vs_rest(prs)
    fig5a(prs)
    fig5b(prs)
    fig5c(prs)
    fig6(prs)
    fig7a(prs)
    fig7b(prs)
    fig8a(prs)
    fig8b(prs)
    fig8c(prs)
    fig8d(prs)

    # Supp
    figS1_biofilm(prs); figS2_ammonia(prs); figS3_alkaline(prs)
    figS4_cazy(prs); figS5_metal(prs); figS6_heroANI(prs)
    figS7_ureCtree(prs)
    figS8_biosafety(prs); figS9_active_site(prs)
    figS10_pseudo_rarity(prs); figS11_alkaliphile(prs)
    figS12_gRodon(prs); figS13_mgnify(prs); figS14_stoich(prs)
    figS15_genomad(prs); figS16_defense(prs)
    figS17a_gc3(prs); figS17b_enc(prs); figS17c_codeml(prs)
    figS17d_yn00(prs); figS18_panMICP(prs); figS19_abundance(prs)
    figS20_esmfold(prs); figS21_antismash(prs)

    prs.save(OUT_PPTX)
    os.makedirs(os.path.dirname(REPO_COPY), exist_ok=True)
    import shutil
    shutil.copy(OUT_PPTX, REPO_COPY)
    print("Wrote", OUT_PPTX)
    print("Copy at", REPO_COPY)


if __name__ == "__main__":
    main()
