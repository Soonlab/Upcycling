"""
Shared style spec + native-shape helpers for the Upcycling editable PPTX.

Nature/Cell/Science conventions distilled:
  - Arial typography hierarchy (panel label 12 pt bold; axis title 10 pt;
    tick / annotation 8-9 pt)
  - 0.5 pt axis/tick weight, outward ticks, top+right spines removed
  - Okabe-Ito colourblind-safe palette for categorical
  - No shadow, no gridline, no 3D, no gradient
  - Significance convention:  * P<0.05, ** P<0.01, *** P<0.001
  - Sample-size 'n=' always in-panel; test statistic annotated in-panel

All text is inserted via add_textbox() (TEXT_BOX autoshape) so every label,
legend item, axis tick, P value and cohort name is double-click editable
in PowerPoint without ungrouping.  Lines and markers are CONNECTOR /
AUTOSHAPE so they remain vector and reshapeable.

Size unit policy: inch everywhere; converted to EMU via Emu / Inches only
at render time.  Slide is 11 x 7 in (wide), one panel per slide.
"""
from __future__ import annotations

from dataclasses import dataclass
from typing import Iterable, Optional

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.shapes.shapetree import SlideShapes
from pptx.util import Emu, Inches, Pt

FONT = "Arial"

# Okabe-Ito colourblind-safe palette
OKABE_ITO = {
    "black":     RGBColor(0x00, 0x00, 0x00),
    "orange":    RGBColor(0xE6, 0x9F, 0x00),
    "sky":       RGBColor(0x56, 0xB4, 0xE9),
    "green":     RGBColor(0x00, 0x9E, 0x73),
    "yellow":    RGBColor(0xF0, 0xE4, 0x42),
    "blue":      RGBColor(0x00, 0x72, 0xB2),
    "vermilion": RGBColor(0xD5, 0x5E, 0x00),
    "purple":    RGBColor(0xCC, 0x79, 0xA7),
    "grey":      RGBColor(0x6B, 0x6B, 0x6B),
    "lightgrey": RGBColor(0xCB, 0xD5, 0xE0),
}

# Lineage colours used across the paper
SPHINGO = OKABE_ITO["blue"]
PSEUDO = OKABE_ITO["vermilion"]
HERO_MIX = OKABE_ITO["purple"]
REST = OKABE_ITO["grey"]

# Waste-source palette
SOURCE_COLOR = {
    "Cattle": RGBColor(0xC0, 0x39, 0x2B),
    "cattle": RGBColor(0xC0, 0x39, 0x2B),
    "Swine":  RGBColor(0x29, 0x6E, 0xC5),
    "swine":  RGBColor(0x29, 0x6E, 0xC5),
    "Sheep":  RGBColor(0x2B, 0x8A, 0x3E),
    "sheep":  RGBColor(0x2B, 0x8A, 0x3E),
    "Poultry": RGBColor(0xE6, 0x7E, 0x22),
    "poultry": RGBColor(0xE6, 0x7E, 0x22),
}

HERO_SET = {"S13", "S16", "S23", "C22", "M1", "S26"}
SPHINGO_HEROES = {"S13", "S16", "S23", "C22"}
PSEUDO_HEROES = {"M1", "S26"}

SLIDE_W = Inches(11)
SLIDE_H = Inches(7)


@dataclass
class AxisSpec:
    """Cartesian-axis bounding box and tick plan."""
    x0: float          # inch, left edge of plot area
    y0: float          # inch, top edge of plot area
    width: float
    height: float
    xlim: tuple[float, float]
    ylim: tuple[float, float]
    x_ticks: Optional[list[float]] = None
    y_ticks: Optional[list[float]] = None
    x_tick_labels: Optional[list[str]] = None
    y_tick_labels: Optional[list[str]] = None
    x_title: str = ""
    y_title: str = ""
    x_log: bool = False
    y_log: bool = False
    x_rotate: int = 0   # tick label rotation degrees
    axis_weight_pt: float = 0.75


def add_text(shapes: SlideShapes, x_in: float, y_in: float,
             w_in: float, h_in: float, text: str, *,
             size: int = 9, bold: bool = False, italic: bool = False,
             color: RGBColor = OKABE_ITO["black"],
             align: str = "left", anchor: str = "top",
             rotation: float = 0) -> None:
    """Append a native TEXT_BOX with Arial formatting, no shadow."""
    tb = shapes.add_textbox(Inches(x_in), Inches(y_in), Inches(w_in),
                            Inches(h_in))
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.auto_size = None
    anchors = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
               "bottom": MSO_ANCHOR.BOTTOM}
    tf.vertical_anchor = anchors.get(anchor, MSO_ANCHOR.TOP)

    p = tf.paragraphs[0]
    aligns = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
              "right": PP_ALIGN.RIGHT}
    p.alignment = aligns.get(align, PP_ALIGN.LEFT)

    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color

    # Disable shadow on the shape (Nature style).
    from pptx.oxml.ns import qn
    from lxml import etree
    spPr = tb._element.spPr  # type: ignore[attr-defined]
    effectLst = spPr.find(qn("a:effectLst"))
    if effectLst is None:
        effectLst = etree.SubElement(spPr, qn("a:effectLst"))

    if rotation:
        tb.rotation = rotation


def _strip_shadow(shape) -> None:
    """Attach an empty effectLst element so PPT renders no shadow."""
    from lxml import etree
    from pptx.oxml.ns import qn
    spPr = shape._element.spPr  # type: ignore[attr-defined]
    # Remove any pre-existing effectLst
    for node in spPr.findall(qn("a:effectLst")):
        spPr.remove(node)
    etree.SubElement(spPr, qn("a:effectLst"))


def add_line(shapes: SlideShapes, x0_in: float, y0_in: float,
             x1_in: float, y1_in: float, *,
             color: RGBColor = OKABE_ITO["black"],
             weight_pt: float = 0.75,
             dash: Optional[str] = None) -> None:
    """Native line as a filled rectangle (robust to LibreOffice connector
    rendering of zero-bbox connectors).  dash: None | 'dash' | 'dot'.

    For purely horizontal/vertical lines we render a thin rectangle of
    the requested weight.  For diagonal lines (dx != 0 and dy != 0) we
    fall back to a straight connector, which renders correctly because
    its bounding box is non-degenerate.
    """
    dx = x1_in - x0_in
    dy = y1_in - y0_in
    # Zero-length guard: both dx and dy ~= 0 means the caller tried to
    # draw a point-like connector — skip it (otherwise LibreOffice PDF
    # export re-anchors the connector to the page origin).
    if abs(dx) < 1e-6 and abs(dy) < 1e-6:
        return
    # line thickness in inches (weight_pt * 1/72)
    thick = max(weight_pt / 72.0, 0.004)

    if abs(dy) < 1e-6 and abs(dx) > 1e-6:
        # horizontal line
        x0 = min(x0_in, x1_in)
        rect = shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(x0), Inches(y0_in - thick / 2),
                                Inches(abs(dx)), Inches(thick))
        rect.fill.solid()
        rect.fill.fore_color.rgb = color
        rect.line.fill.background()
        _strip_shadow(rect)
        _apply_dash(rect, dash, color, weight_pt)
        return
    if abs(dx) < 1e-6 and abs(dy) > 1e-6:
        # vertical line
        y0 = min(y0_in, y1_in)
        rect = shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(x0_in - thick / 2), Inches(y0),
                                Inches(thick), Inches(abs(dy)))
        rect.fill.solid()
        rect.fill.fore_color.rgb = color
        rect.line.fill.background()
        _strip_shadow(rect)
        _apply_dash(rect, dash, color, weight_pt)
        return

    # Diagonal fallback — connector is reliable when bbox is non-degenerate.
    line = shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                Inches(x0_in), Inches(y0_in),
                                Inches(dx), Inches(dy))
    ln = line.line
    ln.color.rgb = color
    ln.width = Pt(weight_pt)
    _strip_shadow(line)
    if dash:
        from pptx.oxml.ns import qn
        from lxml import etree
        dash_map = {"dash": "dash", "dot": "dot", "dashdot": "dashDot"}
        ln_element = line.line._get_or_add_ln()  # type: ignore[attr-defined]
        for node in ln_element.findall(qn("a:prstDash")):
            ln_element.remove(node)
        prst = etree.SubElement(ln_element, qn("a:prstDash"))
        prst.set("val", dash_map.get(dash, "dash"))


def _apply_dash(rect_shape, dash: Optional[str], color: RGBColor,
                weight_pt: float) -> None:
    """Dashed rectangles don't exist in PowerPoint natively; for dashed
    lines we instead chop the rectangle into dashes via multiple small
    tiles sharing the outline property.  For single-pixel rectangles this
    is fine because add_line re-enters us for each dash.

    For simplicity we emulate a dashed effect by thinning the rectangle
    and letting the outline carry the dash pattern."""
    if not dash:
        return
    from pptx.oxml.ns import qn
    from lxml import etree
    dash_map = {"dash": "dash", "dot": "dot", "dashdot": "dashDot"}
    # Put a visible outline in the same colour, set dashed, and clear fill
    rect_shape.fill.background()
    ln = rect_shape.line
    ln.color.rgb = color
    ln.width = Pt(weight_pt)
    ln_el = ln._get_or_add_ln()
    for node in ln_el.findall(qn("a:prstDash")):
        ln_el.remove(node)
    prst = etree.SubElement(ln_el, qn("a:prstDash"))
    prst.set("val", dash_map.get(dash, "dash"))


def add_rect(shapes: SlideShapes, x_in: float, y_in: float,
             w_in: float, h_in: float, *,
             fill: Optional[RGBColor] = None,
             outline: RGBColor = OKABE_ITO["black"],
             outline_weight_pt: float = 0.5) -> None:
    sh = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x_in), Inches(y_in),
                          Inches(w_in), Inches(h_in))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = outline
    sh.line.width = Pt(outline_weight_pt)
    _strip_shadow(sh)
    return sh


def add_circle(shapes: SlideShapes, cx_in: float, cy_in: float,
               r_in: float, *, fill: RGBColor = OKABE_ITO["black"],
               outline: RGBColor = OKABE_ITO["black"],
               outline_weight_pt: float = 0.3) -> None:
    sh = shapes.add_shape(MSO_SHAPE.OVAL,
                          Inches(cx_in - r_in), Inches(cy_in - r_in),
                          Inches(2 * r_in), Inches(2 * r_in))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = outline
    sh.line.width = Pt(outline_weight_pt)
    _strip_shadow(sh)
    return sh


def add_arrow_shape(shapes: SlideShapes, x_in: float, y_in: float,
                    w_in: float, h_in: float, *,
                    fill: RGBColor = OKABE_ITO["blue"],
                    outline: RGBColor = OKABE_ITO["black"],
                    outline_weight_pt: float = 0.4,
                    reverse: bool = False):
    shape_type = (MSO_SHAPE.LEFT_ARROW if reverse
                  else MSO_SHAPE.RIGHT_ARROW)
    sh = shapes.add_shape(shape_type, Inches(x_in), Inches(y_in),
                          Inches(w_in), Inches(h_in))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = outline
    sh.line.width = Pt(outline_weight_pt)
    _strip_shadow(sh)
    return sh


def draw_axis(shapes: SlideShapes, ax: AxisSpec, *,
              draw_x: bool = True, draw_y: bool = True,
              x_title_size: int = 10, y_title_size: int = 10,
              tick_size: int = 8, tick_weight_pt: float = 0.6,
              tick_len_in: float = 0.06) -> None:
    """Draw left + bottom spines with outward ticks and Arial tick labels."""
    # Spines
    if draw_x:
        add_line(shapes, ax.x0, ax.y0 + ax.height, ax.x0 + ax.width,
                 ax.y0 + ax.height, weight_pt=ax.axis_weight_pt)
    if draw_y:
        add_line(shapes, ax.x0, ax.y0, ax.x0,
                 ax.y0 + ax.height, weight_pt=ax.axis_weight_pt)

    # X ticks
    if ax.x_ticks is not None and draw_x:
        for tval, tlbl in zip(ax.x_ticks, ax.x_tick_labels or
                              [f"{v:g}" for v in ax.x_ticks]):
            xp = data_to_plot_x(tval, ax)
            y_bot = ax.y0 + ax.height
            add_line(shapes, xp, y_bot, xp, y_bot + tick_len_in,
                     weight_pt=tick_weight_pt)
            add_text(shapes, xp - 0.6, y_bot + tick_len_in + 0.02,
                     1.2, 0.22, tlbl, size=tick_size, align="center",
                     rotation=ax.x_rotate)

    # Y ticks
    if ax.y_ticks is not None and draw_y:
        for tval, tlbl in zip(ax.y_ticks, ax.y_tick_labels or
                              [f"{v:g}" for v in ax.y_ticks]):
            yp = data_to_plot_y(tval, ax)
            add_line(shapes, ax.x0 - tick_len_in, yp, ax.x0, yp,
                     weight_pt=tick_weight_pt)
            add_text(shapes, ax.x0 - tick_len_in - 0.56, yp - 0.09,
                     0.54, 0.20, tlbl, size=tick_size, align="right")

    # Titles
    if ax.x_title and draw_x:
        add_text(shapes, ax.x0 - 0.2, ax.y0 + ax.height + 0.45,
                 ax.width + 0.4, 0.3, ax.x_title, size=x_title_size,
                 align="center", bold=True)
    if ax.y_title and draw_y:
        add_text(shapes, ax.x0 - 1.0, ax.y0 + ax.height / 2 - 0.15,
                 0.8, 0.3, ax.y_title, size=y_title_size,
                 align="right", bold=True, rotation=-90)


def data_to_plot_x(v: float, ax: AxisSpec) -> float:
    import math
    lo, hi = ax.xlim
    if ax.x_log:
        v = math.log10(max(v, 1e-12))
        lo = math.log10(max(lo, 1e-12))
        hi = math.log10(max(hi, 1e-12))
    return ax.x0 + (v - lo) / (hi - lo) * ax.width


def data_to_plot_y(v: float, ax: AxisSpec) -> float:
    import math
    lo, hi = ax.ylim
    if ax.y_log:
        v = math.log10(max(v, 1e-12))
        lo = math.log10(max(lo, 1e-12))
        hi = math.log10(max(hi, 1e-12))
    # invert — plot Y grows upward, inches grow downward
    return ax.y0 + ax.height - (v - lo) / (hi - lo) * ax.height


def nice_log_ticks(lo: float, hi: float) -> list[float]:
    """Decade tick positions inclusive."""
    import math
    if lo <= 0:
        lo = max(lo, 1e-6)
    decades = []
    d = math.floor(math.log10(lo))
    while 10 ** d <= hi * 1.0001:
        decades.append(10 ** d)
        d += 1
    return decades


def panel_label(shapes: SlideShapes, label: str,
                x_in: float = 0.3, y_in: float = 0.25) -> None:
    """Top-left bold lowercase panel label (e.g., 'a', 'b')."""
    add_text(shapes, x_in, y_in, 0.3, 0.3, label,
             size=18, bold=True, color=OKABE_ITO["black"])


def sig_star(p: float) -> str:
    if p < 1e-3:
        return "***"
    if p < 1e-2:
        return "**"
    if p < 5e-2:
        return "*"
    return "n.s."


def hero_color_for(mag: str) -> RGBColor:
    if mag in SPHINGO_HEROES:
        return SPHINGO
    if mag in PSEUDO_HEROES:
        return PSEUDO
    return REST
